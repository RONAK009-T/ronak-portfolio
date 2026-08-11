import os
import sys
import subprocess
import signal
from django.shortcuts import render, get_object_or_404, redirect
from django.http import JsonResponse, HttpResponse
from .models import Project, Skill, Experience, Education, Service, SocialLink

# Python Excel & PowerPoint generation libraries
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Keep track of active server subprocesses locally (key: project_id -> value: subprocess.Popen)
running_servers = {}

def index_view(request):
    projects = Project.objects.all().order_by('-featured', '-created_date')
    skills = Skill.objects.all()
    experience = Experience.objects.all().order_by('order', '-id')
    education = Education.objects.all().order_by('order', '-id')
    services = Service.objects.all().order_by('order')
    social_links = SocialLink.objects.all().order_by('order')

    # Categorize skills for grid layouts
    frontend_skills = skills.filter(category='Frontend')
    backend_skills = skills.filter(category='Backend')
    ai_skills = skills.filter(category='AI/Data Science')
    tool_skills = skills.filter(category='Tools')

    context = {
        'projects': projects,
        'frontend_skills': frontend_skills,
        'backend_skills': backend_skills,
        'ai_skills': ai_skills,
        'tool_skills': tool_skills,
        'experience': experience,
        'education': education,
        'services': services,
        'social_links': social_links,
    }
    return render(request, 'portfolio/index.html', context)


def project_detail_view(request, slug):
    project = get_object_or_404(Project, slug=slug)
    # Get other projects for recommendation carousel
    recent_projects = Project.objects.exclude(id=project.id)[:3]
    return render(request, 'portfolio/project_detail.html', {
        'project': project,
        'recent_projects': recent_projects
    })


def start_project_api(request, slug):
    """
    API view that launches a local project (Django, Streamlit, React, or Static HTML)
    in a background subprocess based on the files present in its directory.
    """
    project = get_object_or_404(Project, slug=slug)
    
    if not project.local_port:
        if project.live_url and project.live_url.startswith('file:///'):
            return JsonResponse({"status": "static", "url": project.live_url})
        return JsonResponse({"status": "error", "message": "This project is not configured for local execution."})

    port = project.local_port
    path = project.local_path

    # If the process is already running, return its address
    if project.id in running_servers:
        proc = running_servers[project.id]
        if proc.poll() is None: # Still running
            return JsonResponse({"status": "running", "url": f"http://127.0.0.1:{port}"})

    if not path or not os.path.exists(path):
        return JsonResponse({
            "status": "error", 
            "message": f"Project folder path not found on server directory: {path}"
        })

    python_exe = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..', 'venv', 'Scripts', 'python.exe'))
    if not os.path.exists(python_exe):
        python_exe = "python" # Fallback

    try:
        # 1. Django Project
        if os.path.exists(os.path.join(path, 'manage.py')):
            proc = subprocess.Popen(
                [python_exe, 'manage.py', 'runserver', f'127.0.0.1:{port}'],
                cwd=path,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                creationflags=subprocess.CREATE_NEW_PROCESS_GROUP if sys.platform == 'win32' else 0
            )
        # 2. Streamlit Project
        elif os.path.exists(os.path.join(path, 'app.py')) and not os.path.exists(os.path.join(path, 'manage.py')):
            streamlit_exe = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..', 'venv', 'Scripts', 'streamlit.exe'))
            if not os.path.exists(streamlit_exe):
                streamlit_exe = "streamlit"
            proc = subprocess.Popen(
                [streamlit_exe, 'run', 'app.py', '--server.port', str(port)],
                cwd=path,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                creationflags=subprocess.CREATE_NEW_PROCESS_GROUP if sys.platform == 'win32' else 0
            )
        # 3. React / Node Project
        elif os.path.exists(os.path.join(path, 'package.json')):
            # Run Vite server directly inside the project directory
            proc = subprocess.Popen(
                ['cmd.exe', '/c', f'npx vite --port {port} --host 127.0.0.1'],
                cwd=path,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                creationflags=subprocess.CREATE_NEW_PROCESS_GROUP if sys.platform == 'win32' else 0
            )
        # 4. Static HTML Project
        elif os.path.exists(os.path.join(path, 'index.html')) or os.path.exists(os.path.join(path, 'home.html')):
            proc = subprocess.Popen(
                [python_exe, '-m', 'http.server', str(port)],
                cwd=path,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                creationflags=subprocess.CREATE_NEW_PROCESS_GROUP if sys.platform == 'win32' else 0
            )
        else:
            return JsonResponse({
                "status": "error", 
                "message": "Unable to determine project structure (no manage.py, app.py, package.json, or index.html found)."
            })

        running_servers[project.id] = proc
        return JsonResponse({"status": "started", "url": f"http://127.0.0.1:{port}"})
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)})



def stop_project_api(request, slug):
    """
    API view that terminates a running background process.
    """
    project = get_object_or_404(Project, slug=slug)
    if project.id in running_servers:
        proc = running_servers[project.id]
        try:
            if sys.platform == 'win32':
                proc.send_signal(signal.CTRL_BREAK_EVENT)
            else:
                proc.terminate()
            proc.wait(timeout=3)
        except Exception:
            proc.kill()
        
        del running_servers[project.id]
        return JsonResponse({"status": "stopped"})
    return JsonResponse({"status": "not_running"})


def generate_excel_api(request):
    """
    B.Com Analysis Skill Integration: Generates a stylized excel sheet 
    detailing a mock financial forecast using openpyxl.
    """
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Financial Forecast"

    # Style definitions (Black + Neon Green theme)
    font_title = Font(name='Courier New', size=16, bold=True, color='00FF88')
    font_header = Font(name='Segoe UI', size=11, bold=True, color='FFFFFF')
    font_cell = Font(name='Segoe UI', size=10, color='D3D3D3')
    fill_header = PatternFill(start_color='101010', end_color='101010', fill_type='solid')
    fill_bg = PatternFill(start_color='000000', end_color='000000', fill_type='solid')
    align_center = Alignment(horizontal='center', vertical='center')

    # Set background style for visual check
    for row in range(1, 15):
        for col in range(1, 7):
            cell = ws.cell(row=row, column=col)
            cell.fill = fill_bg

    ws.merge_cells('A2:F2')
    ws['A2'] = "RONAK // B.COM FINANCIAL SUITE - PROJECTIONS"
    ws['A2'].font = font_title
    ws['A2'].alignment = align_center

    headers = ["MONTH", "REVENUE ($)", "COSTS ($)", "GROSS MARGIN ($)", "TAX (15%)", "NET PROFIT ($)"]
    for col_idx, header in enumerate(headers, 1):
        cell = ws.cell(row=4, column=col_idx, value=header)
        cell.font = font_header
        cell.fill = fill_header
        cell.alignment = align_center

    months = [("Jan", 12000, 5000), ("Feb", 15000, 6000), ("Mar", 18000, 7200), ("Apr", 22000, 8500)]
    for row_idx, (month, rev, cost) in enumerate(months, 5):
        ws.cell(row=row_idx, column=1, value=month).font = font_cell
        ws.cell(row=row_idx, column=2, value=rev).font = font_cell
        ws.cell(row=row_idx, column=3, value=cost).font = font_cell
        # Excel Formulas
        ws.cell(row=row_idx, column=4, value=f"=B{row_idx}-C{row_idx}").font = font_cell
        ws.cell(row=row_idx, column=5, value=f"=D{row_idx}*0.15").font = font_cell
        ws.cell(row=row_idx, column=6, value=f"=D{row_idx}-E{row_idx}").font = font_cell

    # Auto-adjust column widths
    for col in ws.columns:
        max_len = max(len(str(cell.value or '')) for cell in col)
        ws.column_dimensions[col[0].column_letter].width = max(max_len + 3, 12)

    response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
    response['Content-Disposition'] = 'attachment; filename=Ronak_Financial_Analysis.xlsx'
    wb.save(response)
    return response


def generate_ppt_api(request):
    """
    B.Com presentation Skill Integration: Generates a slide summary of Ronak's profile.
    """
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Ronak - Digital Portfolio"
    subtitle.text = "B.Com Student (Final Year) & Full Stack Python/Django/React Developer"

    # Slide 2: Technical Constellation
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Core Competencies"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Python Development & Django Framework Integration"
    p1 = tf.add_paragraph()
    p1.text = "React Frontend & Responsive Grid Designs"
    p2 = tf.add_paragraph()
    p2.text = "Business Intelligence: Advanced Microsoft Excel & PowerPoint Presentations"
    p3 = tf.add_paragraph()
    p3.text = "Financial Modeling & Commercial Data Analytics"

    response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = 'attachment; filename=Ronak_Portfolio_Presentation.pptx'
    prs.save(response)
    return response
